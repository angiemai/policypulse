import fitz 
import docx
import pytesseract
from pdf2image import convert_from_path
from pptx import Presentation
from odf import text, teletype
from odf.opendocument import load
import os
from pathlib import Path
from tqdm.notebook import tqdm
import datetime
from pinecone import Pinecone

# Function to extract text from DOCX
def extract_text_from_docx(file_path):
    doc = docx.Document(file_path)
    text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
    return text

# Function to extract text from PPTX
def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

# Function to extract text from ODP
def extract_text_from_odp(file_path):
    doc = load(file_path)
    text_elements = doc.getElementsByType(text.P)
    return "\n".join([teletype.extractText(element) for element in text_elements])

# Function to check if a PDF is scanned and extract text accordingly
def extract_text_from_pdf(file_path):
    # Open the PDF
    doc = fitz.open(file_path)
    
    # Check if the PDF has text
    text = ""
    text_found = False
    
    for page_num in range(len(doc)):
        page = doc[page_num]
        page_text = page.get_text()
        
        # If page has more than 10 characters, consider it a text PDF
        if len(page_text.strip()) > 10:
            text_found = True
            text += page_text + "\n"
    
    # If no significant text found, it's likely a scanned PDF
    if not text_found:
        print(f"PDF appears to be scanned, applying OCR: {file_path}")
        return extract_text_from_scanned_pdf(file_path)
    
    return text

# Function to extract text from scanned PDFs using OCR
def extract_text_from_scanned_pdf(file_path):
    # Convert PDF to images
    images = convert_from_path(file_path)
    
    # Apply OCR to each image
    text = []
    for i, image in enumerate(images):
        text.append(pytesseract.image_to_string(image))
    
    return "\n".join(text)

# Main function to extract text based on file extension
def extract_text(file_path):
    file_path = Path(file_path)
    extension = file_path.suffix.lower()
    
    try:
        if extension == '.pdf':
            return extract_text_from_pdf(file_path)
        elif extension == '.docx':
            return extract_text_from_docx(file_path)
        elif extension == '.pptx':
            return extract_text_from_pptx(file_path)
        elif extension == '.odp':
            return extract_text_from_odp(file_path)
        else:
            print(f"Unsupported file type: {extension}")
            return None
    except Exception as e:
        print(f"Error extracting text from {file_path}: {e}")
        return None


# Function to process a file: extract text and split into chunks
def process_file(file_path, text_splitter):
    """
    Process a single file: extract text and split into chunks
    
    Parameters:
    file_path (str or Path): Path to the file to process
    text_splitter: Initialized text splitter object
    
    Returns:
    list: List of document dictionaries with text and metadata
    """
    # Extract text
    text = extract_text(file_path)
    if not text:
        return []
    
    # Split text into chunks
    chunks = text_splitter.split_text(text)
    
    # Create metadata
    file_path = Path(file_path)
    metadata = {
        "filename": file_path.name,
        "file_path": str(file_path),
        "file_type": file_path.suffix.lower(),
        # Format timestamps as ISO format strings
        "created_at": datetime.datetime.fromtimestamp(os.path.getctime(file_path)).isoformat(),
        "modified_at": datetime.datetime.fromtimestamp(os.path.getmtime(file_path)).isoformat(),
    }
    
    # Create document objects
    documents = []
    for i, chunk in enumerate(chunks):
        # Create a unique chunk metadata
        chunk_metadata = metadata.copy()
        chunk_metadata["chunk_id"] = i
        chunk_metadata["chunk_text"] = chunk[:100] + "..."  # Preview of the chunk
        
        documents.append({
            "id": f"{file_path.stem}_chunk_{i}",
            "text": chunk,  # This is now the raw text, not vectors
            "metadata": chunk_metadata
        })
    
    return documents

# Main function to process all supported documents in a directory
def process_directory(directory_path, text_splitter):
    """
    Process all supported documents in a directory
    
    Parameters:
    directory_path (str or Path): Path to directory containing documents
    text_splitter: Initialized text splitter object
    
    Returns:
    list: Combined list of document objects from all processed files
    """
    # Supported extensions
    supported_extensions = ['.pdf', '.docx', '.pptx', '.odp']
    
    # Get all files
    directory = Path(directory_path)
    all_files = [p for p in directory.glob('**/*') if p.is_file() and p.suffix.lower() in supported_extensions]
    
    print(f"Found {len(all_files)} supported documents to process")
    
    all_documents = []
    
    # Process each file
    for file_path in tqdm(all_files, desc="Processing files"):
        print(f"Processing: {file_path}")
        file_documents = process_file(file_path, text_splitter)
        all_documents.extend(file_documents)
        print(f"Created {len(file_documents)} chunks for {file_path}")
    
    print(f"Total document chunks created: {len(all_documents)}")
    return all_documents


# Function to upload embeddings to Pinecone
def upload_to_pinecone(documents, index_name, api_key):
    """
    Upload documents to Pinecone using text-based upsert
    
    Parameters:
    documents (list): List of document objects with text and metadata
    index_name (str): Name of Pinecone index
    api_key (str): Pinecone API key
    """
    #from pinecone import Pinecone
    import json
    
    # Initialize Pinecone
    pc = Pinecone(api_key=api_key)
    
    # Get the index
    index = pc.Index(index_name)
    
    # Calculate number of batches
    batch_size = 64  # Pinecone appears to have limit of 96
    num_batches = (len(documents) - 1) // batch_size + 1
    
    # Upsert documents (Pinecone will handle the embedding generation)
    for i in tqdm(range(0, len(documents), batch_size), desc="Uploading to Pinecone", total=num_batches):
        batch = documents[i:i+batch_size]
        
        # Format for text-based upsert with metadata as a single string
        upsert_batch = []
        for doc in batch:
            # Create a simplified metadata dictionary
            metadata_dict = {
                "filename": doc["metadata"]["filename"],
                "file_type": doc["metadata"]["file_type"],
                "chunk_id": str(doc["metadata"]["chunk_id"]),
                "preview": doc["metadata"]["chunk_text"]
            }
            
            # Convert metadata to a JSON string with triple quotes
            metadata_str = json.dumps(metadata_dict)
            
            # Create record with metadata as a single string
            record = {
                "id": doc["id"],
                "text": doc["text"],  # The raw text
                "metadata": metadata_str  # Metadata as a serialized JSON string
            }
            
            upsert_batch.append(record)
        
        if i == 0:  # Print sample of first batch only
            print("Sample record format:", upsert_batch[0])
            
        # Upsert to Pinecone
        index.upsert_records("", upsert_batch)
    
  